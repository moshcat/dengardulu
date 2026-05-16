"""Generate DengarDulu pitch-deck.pptx using Mastercard-inspired design system.

Reads /pitch-deck.md content and applies the DESIGN.MD visual language:
  - Canvas Cream #F3F0EE page canvas (never pure white)
  - Ink Black #141413 primary text + CTA pills
  - Signal Light Orange #F37338 accents / orbital cues
  - Signal Orange #CF4500 danger / signal
  - Pill / stadium shapes (999px-equivalent radius) + 40px hero radius
  - Sofia Sans typography (open-source MarkForMC stand-in)
  - Eyebrow labels with a dot + uppercase +4% tracking
  - Editorial whitespace, one idea per slide, asymmetric placement

Run:  python3 scripts/build-pitch-deck.py
Output: pitch-deck.pptx in project root
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ---------- Design tokens ----------
CANVAS_CREAM = RGBColor(0xF3, 0xF0, 0xEE)
LIFTED_CREAM = RGBColor(0xFC, 0xFB, 0xFA)
INK_BLACK = RGBColor(0x14, 0x14, 0x13)
CHARCOAL = RGBColor(0x26, 0x26, 0x27)
SLATE_GRAY = RGBColor(0x69, 0x69, 0x69)
SIGNAL_ORANGE = RGBColor(0xCF, 0x45, 0x00)
SIGNAL_LIGHT = RGBColor(0xF3, 0x73, 0x38)
CLAY_BROWN = RGBColor(0x9A, 0x3A, 0x0A)
DUST_TAUPE = RGBColor(0xD1, 0xCD, 0xC7)
SUCCESS_GREEN = RGBColor(0x0A, 0x7A, 0x3D)
DANGER_RED = RGBColor(0xEB, 0x00, 0x1B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GHOST_CREAM = RGBColor(0xE8, 0xE2, 0xDA)

HEAD_FONT = "Sofia Sans"
BODY_FONT = "Sofia Sans"
MONO_FONT = "Inter"


# ---------- Helpers ----------
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0, shape=MSO_SHAPE.RECTANGLE, adj=None):
    shp = slide.shapes.add_shape(shape, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    if adj is not None and hasattr(shp, "adjustments"):
        try:
            shp.adjustments[0] = adj
        except Exception:
            pass
    return shp


def add_pill(slide, x, y, w, h, fill=INK_BLACK, line=None):
    """Full stadium/pill shape."""
    shp = add_rect(slide, x, y, w, h, fill=fill, line=line,
                   line_w=1.5 if line else 0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.5)
    return shp


def add_rounded(slide, x, y, w, h, fill=None, line=None, radius_in=0.4):
    """Rounded rectangle with a specific corner radius (approx, expressed as fraction)."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.5)
    short = min(w, h)
    adj = min(0.5, max(0.0, Emu(Inches(radius_in)) / short))
    try:
        shp.adjustments[0] = adj
    except Exception:
        pass
    return shp


def add_text(slide, x, y, w, h, text, *, size=16, bold=False, italic=False, color=INK_BLACK,
             font=BODY_FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             tracking=None, line_spacing=None):
    """Add a styled text box. `tracking` in hundredths-of-a-point (100 = 1pt)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    tb.text_frame.margin_left = 0
    tb.text_frame.margin_right = 0
    tb.text_frame.margin_top = 0
    tb.text_frame.margin_bottom = 0
    tb.text_frame.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else text
    first = True
    for ln in lines:
        p = tb.text_frame.paragraphs[0] if first else tb.text_frame.add_paragraph()
        p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = ln
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        if tracking is not None:
            _set_spacing(run, tracking)
        first = False
    return tb


def add_rich(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=None, paragraph_space_after=None):
    """Add a text box with multiple runs / paragraphs.

    `runs` is a list. Each item is either:
      - dict: a single run on the current paragraph
      - list of dicts: multiple runs making up one paragraph
      - 'BREAK': insert paragraph break
    Run dict keys: text, size, bold, color, font, tracking, align
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    first_para = True
    current_para = None

    def new_para(override_align=None):
        nonlocal first_para, current_para
        if first_para:
            current_para = tf.paragraphs[0]
            first_para = False
        else:
            current_para = tf.add_paragraph()
        current_para.alignment = override_align if override_align is not None else align
        if line_spacing is not None:
            current_para.line_spacing = line_spacing
        if paragraph_space_after is not None:
            current_para.space_after = Pt(paragraph_space_after)
        return current_para

    def add_run_to_para(para, spec):
        r = para.add_run()
        r.text = spec.get("text", "")
        r.font.name = spec.get("font", BODY_FONT)
        r.font.size = Pt(spec.get("size", 16))
        r.font.bold = spec.get("bold", False)
        r.font.italic = spec.get("italic", False)
        r.font.color.rgb = spec.get("color", INK_BLACK)
        if spec.get("tracking") is not None:
            _set_spacing(r, spec["tracking"])

    for item in runs:
        if item == "BREAK":
            new_para()
            continue
        if isinstance(item, dict):
            para_align = item.get("align")
            para = new_para(para_align)
            add_run_to_para(para, item)
        elif isinstance(item, list):
            para_align = item[0].get("align") if item else None
            para = new_para(para_align)
            for spec in item:
                add_run_to_para(para, spec)
    return tb


def _set_spacing(run, hundredths_pt):
    """Set character letter-spacing (OOXML spc attr)."""
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(hundredths_pt)))


def add_eyebrow(slide, x, y, text, color=SIGNAL_LIGHT, text_color=SLATE_GRAY):
    """Eyebrow label: • LABEL (uppercase, +4% tracking, weight 700, 14px)."""
    # dot circle
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Inches(0.05), Inches(0.1), Inches(0.1))
    dot.shadow.inherit = False
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    add_text(slide, x + Inches(0.18), y, Inches(8), Inches(0.3),
             text.upper(), size=10, bold=True, color=text_color,
             font=HEAD_FONT, tracking=56)  # +4% on 14pt ≈ 56 hundredths


def add_ink_pill_button(slide, x, y, w, h, label, *, fill=INK_BLACK, text_color=CANVAS_CREAM, size=12):
    add_pill(slide, x, y, w, h, fill=fill)
    add_text(slide, x, y, w, h, label, size=size, bold=True, color=text_color,
             font=HEAD_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_orbit_arc(slide, x, y, w, h, color=SIGNAL_LIGHT, weight_pt=1.25):
    """Thin decorative arc; approximated by a quarter-circle arc shape."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ARC, x, y, w, h)
    shp.shadow.inherit = False
    shp.fill.background()
    shp.line.color.rgb = color
    shp.line.width = Pt(weight_pt)


def add_page_number(slide, idx, total):
    add_text(slide, Inches(12.55), Inches(7.08), Inches(0.7), Inches(0.3),
             f"{idx:02d} / {total:02d}", size=9, color=SLATE_GRAY,
             font=HEAD_FONT, align=PP_ALIGN.RIGHT, tracking=40)


def add_footer_mark(slide):
    # tiny brand mark bottom-left
    add_text(slide, Inches(0.5), Inches(7.08), Inches(4), Inches(0.3),
             "DENGARDULU  ·  LISTEN FIRST. ANSWER WISELY.",
             size=9, color=SLATE_GRAY, font=HEAD_FONT, tracking=40, bold=True)


def add_soundwave(slide, cx, cy, width, bars=9, color=INK_BLACK):
    """Decorative micro-soundwave mark (vertical bars with varying heights)."""
    import math
    total_w = width
    gap = Emu(total_w / (bars * 2.2))
    bar_w = Emu(total_w / (bars * 2.0))
    start_x = cx - Emu(total_w / 2)
    heights = [0.35, 0.6, 0.85, 1.0, 0.75, 1.0, 0.55, 0.75, 0.4]
    for i in range(bars):
        h = Inches(heights[i % len(heights)] * 0.4)
        bx = start_x + i * (bar_w + gap)
        by = cy - Emu(h / 2)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, bar_w, h)
        bar.shadow.inherit = False
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        try:
            bar.adjustments[0] = 0.5
        except Exception:
            pass


# ---------- Deck setup ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

TOTAL_SLIDES = 15


def new_slide():
    s = prs.slides.add_slide(BLANK)
    set_slide_bg(s, CANVAS_CREAM)
    return s


# ---------- SLIDE 1 — TITLE ----------
s = new_slide()

# Ghost watermark behind
add_text(s, Inches(-0.5), Inches(0.6), Inches(14), Inches(3),
         "dengardulu", size=220, bold=True, color=GHOST_CREAM,
         font=HEAD_FONT, align=PP_ALIGN.CENTER, tracking=-440)

# Eyebrow
add_eyebrow(s, Inches(6.0), Inches(1.6), "• PROJECT 2030 · MYAI FUTURE HACKATHON",
            color=SIGNAL_LIGHT, text_color=SLATE_GRAY)

# Decorative soundwave
add_soundwave(s, Inches(6.666), Inches(2.5), Inches(1.8), bars=9, color=INK_BLACK)

# Main wordmark
add_text(s, Inches(1), Inches(2.9), Inches(11.333), Inches(1.8),
         "DengarDulu", size=104, bold=True, color=INK_BLACK,
         font=HEAD_FONT, align=PP_ALIGN.CENTER, tracking=-208,
         line_spacing=1.0)

# Tagline — Signal Light
add_text(s, Inches(1), Inches(4.85), Inches(11.333), Inches(0.7),
         "Listen First. Answer Wisely.", size=28, bold=False, color=SIGNAL_LIGHT,
         font=HEAD_FONT, align=PP_ALIGN.CENTER, tracking=-56)

# Subtitle
add_text(s, Inches(1), Inches(5.6), Inches(11.333), Inches(0.5),
         "AI voice-scam safety copilot for Malaysia",
         size=16, color=CHARCOAL, font=BODY_FONT, align=PP_ALIGN.CENTER)

# Track pill
add_pill(s, Inches(4.4), Inches(6.25), Inches(4.5), Inches(0.4), fill=INK_BLACK)
add_text(s, Inches(4.4), Inches(6.25), Inches(4.5), Inches(0.4),
         "TRACK 5  ·  SECURE DIGITAL (FINTECH & SECURITY)",
         size=10, bold=True, color=CANVAS_CREAM, font=HEAD_FONT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, tracking=80)

# Team footer
add_text(s, Inches(1), Inches(6.85), Inches(11.333), Inches(0.4),
         "GDG On Campus UTM  ·  April 2026",
         size=11, color=SLATE_GRAY, font=HEAD_FONT, align=PP_ALIGN.CENTER)

add_page_number(s, 1, TOTAL_SLIDES)


# ---------- SLIDE 2 — PROBLEM STATS ----------
s = new_slide()
add_eyebrow(s, Inches(0.7), Inches(0.55), "• THE PROBLEM")
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(0.9),
         "The scale of the problem", size=46, bold=True, color=INK_BLACK,
         font=HEAD_FONT, tracking=-92)

# Three big stat cards (pills with number + label)
cards = [
    ("454+", "AI voice-cloning cases investigated\nby CCID since 2024"),
    ("RM 2.72M", "Documented losses,\nand climbing"),
    ("12,110", "Scams in Q1 2025 alone.\nRM 573.7M stolen."),
]
card_w = Inches(3.8)
card_h = Inches(3.8)
gap = Inches(0.35)
total_w = card_w * 3 + gap * 2
start_x = (prs.slide_width - total_w) / 2
card_y = Inches(2.3)

for i, (num, label) in enumerate(cards):
    cx = start_x + i * (card_w + gap)
    add_rounded(s, cx, card_y, card_w, card_h, fill=LIFTED_CREAM, radius_in=0.55)
    # Eyebrow inside card
    add_eyebrow(s, cx + Inches(0.45), card_y + Inches(0.45),
                "• MALAYSIA  ·  2024–2025", text_color=SLATE_GRAY)
    # Big number
    add_text(s, cx, card_y + Inches(1.0), card_w, Inches(1.4), num,
             size=64 if len(num) <= 6 else 56, bold=True, color=SIGNAL_LIGHT if i == 1 else INK_BLACK,
             font=HEAD_FONT, align=PP_ALIGN.CENTER, tracking=-130)
    # Label
    add_text(s, cx + Inches(0.4), card_y + Inches(2.6), card_w - Inches(0.8), Inches(1.1),
             label, size=14, color=CHARCOAL, font=BODY_FONT,
             align=PP_ALIGN.CENTER, line_spacing=1.35)

# Source line
add_text(s, Inches(0.7), Inches(6.65), Inches(12), Inches(0.4),
         "Source: CCID Malaysia  ·  Bank Negara Malaysia  ·  Malay Mail 2025",
         size=10, color=SLATE_GRAY, font=HEAD_FONT, align=PP_ALIGN.CENTER,
         tracking=40)

add_footer_mark(s)
add_page_number(s, 2, TOTAL_SLIDES)


# ---------- SLIDE 3 — REAL MALAYSIAN CASE ----------
s = new_slide()
add_eyebrow(s, Inches(0.7), Inches(0.5), "• A REAL MALAYSIAN CASE")
add_text(s, Inches(0.7), Inches(0.85), Inches(8.5), Inches(2.2),
         "A Selangor woman\nlost RM 5,000\nlast year.",
         size=44, bold=True, color=INK_BLACK, font=HEAD_FONT,
         tracking=-88, line_spacing=1.05)

# Narrative block left
story_runs = [
    [
        {"text": "Her ", "size": 16, "color": CHARCOAL},
        {"text": "boss's voice", "size": 16, "bold": True, "color": INK_BLACK},
        {"text": " — cloned from a public TikTok —", "size": 16, "color": CHARCOAL},
    ],
    [
        {"text": "asked her via WhatsApp to buy Touch 'n Go PINs.", "size": 16, "color": CHARCOAL},
    ],
    "BREAK",
    [
        {"text": "She recognized the voice. She trusted it.", "size": 16, "color": CHARCOAL, "italic": True},
    ],
    "BREAK",
    [
        {"text": "By the time she realized it wasn't her boss,", "size": 16, "color": CHARCOAL},
    ],
    [
        {"text": "the PINs were already redeemed.", "size": 16, "bold": True, "color": SIGNAL_ORANGE},
    ],
]
add_rich(s, Inches(0.7), Inches(3.55), Inches(7.6), Inches(2.8),
         story_runs, line_spacing=1.35, paragraph_space_after=3)

# Kicker line
add_text(s, Inches(0.7), Inches(6.55), Inches(8.5), Inches(0.4),
         "This is a real case. And there are 453 more like it.",
         size=12, color=SLATE_GRAY, font=BODY_FONT, italic=True)

# WhatsApp mock card — right side
mock_x = Inches(9.3)
mock_y = Inches(1.1)
mock_w = Inches(3.4)
mock_h = Inches(5.6)
add_rounded(s, mock_x, mock_y, mock_w, mock_h, fill=INK_BLACK, radius_in=0.55)

# Phone status pill
add_pill(s, mock_x + Inches(1.35), mock_y + Inches(0.3), Inches(0.7), Inches(0.18),
         fill=CHARCOAL)

# Chat header
add_text(s, mock_x + Inches(0.3), mock_y + Inches(0.7), mock_w - Inches(0.6), Inches(0.3),
         "Boss",
         size=14, bold=True, color=CANVAS_CREAM, font=HEAD_FONT)
add_text(s, mock_x + Inches(0.3), mock_y + Inches(1.0), mock_w - Inches(0.6), Inches(0.3),
         "online",
         size=9, color=SIGNAL_LIGHT, font=BODY_FONT)

# Voice bubble
bubble_y = mock_y + Inches(1.6)
add_rounded(s, mock_x + Inches(0.3), bubble_y, mock_w - Inches(0.6), Inches(1.0),
            fill=LIFTED_CREAM, radius_in=0.3)
# play button
play = s.shapes.add_shape(MSO_SHAPE.OVAL, mock_x + Inches(0.5), bubble_y + Inches(0.25),
                          Inches(0.5), Inches(0.5))
play.shadow.inherit = False
play.fill.solid()
play.fill.fore_color.rgb = SIGNAL_LIGHT
play.line.fill.background()
add_text(s, mock_x + Inches(0.5), bubble_y + Inches(0.25), Inches(0.5), Inches(0.5),
         "▶", size=16, color=WHITE, font=HEAD_FONT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# waveform simulated
add_text(s, mock_x + Inches(1.1), bubble_y + Inches(0.3), Inches(1.8), Inches(0.4),
         "▁▃▅▇▆▄▅▇▃▂▄▆▇▅▃", size=14, color=SLATE_GRAY, font=MONO_FONT)
add_text(s, mock_x + Inches(1.1), bubble_y + Inches(0.62), Inches(1.6), Inches(0.3),
         "0:23",
         size=10, color=SLATE_GRAY, font=BODY_FONT)

# Urgent-looking text bubble
bubble2_y = bubble_y + Inches(1.3)
add_rounded(s, mock_x + Inches(0.3), bubble2_y, mock_w - Inches(0.6), Inches(0.9),
            fill=LIFTED_CREAM, radius_in=0.3)
add_text(s, mock_x + Inches(0.5), bubble2_y + Inches(0.15), mock_w - Inches(1.0), Inches(0.7),
         "Cepat lah, urgent ni.\nBeli TNG ePIN, hantar balik.",
         size=11, color=INK_BLACK, font=BODY_FONT, line_spacing=1.3)

# Danger overlay
warn_y = bubble2_y + Inches(1.15)
add_pill(s, mock_x + Inches(0.3), warn_y, mock_w - Inches(0.6), Inches(0.5),
         fill=SIGNAL_ORANGE)
add_text(s, mock_x + Inches(0.3), warn_y, mock_w - Inches(0.6), Inches(0.5),
         "⚠  HIGH RISK — CLONED VOICE",
         size=11, bold=True, color=WHITE, font=HEAD_FONT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, tracking=40)

# Small caption
add_text(s, mock_x, mock_y + mock_h + Inches(0.05), mock_w, Inches(0.3),
         "Reconstruction of a real WhatsApp message",
         size=9, color=SLATE_GRAY, font=BODY_FONT, align=PP_ALIGN.CENTER,
         tracking=20)

add_footer_mark(s)
add_page_number(s, 3, TOTAL_SLIDES)


# ---------- SLIDE 4 — WHY EXISTING SOLUTIONS FAIL ----------
s = new_slide()
add_eyebrow(s, Inches(0.7), Inches(0.55), "• THE GAP")
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(1.0),
         "Why existing defenses break down",
         size=42, bold=True, color=INK_BLACK, font=HEAD_FONT, tracking=-80)

col_w = Inches(5.8)
col_y = Inches(2.3)
col_h = Inches(4.6)

# LEFT card — Typical approach (grey, X'd)
left_x = Inches(0.7)
add_rounded(s, left_x, col_y, col_w, col_h, fill=LIFTED_CREAM, radius_in=0.5)
add_eyebrow(s, left_x + Inches(0.5), col_y + Inches(0.5),
            "• TYPICAL APPROACH", color=DUST_TAUPE, text_color=SLATE_GRAY)
add_text(s, left_x + Inches(0.5), col_y + Inches(0.95),
         col_w - Inches(1), Inches(0.8),
         "Binary deepfake detection",
         size=24, bold=True, color=SLATE_GRAY, font=HEAD_FONT, tracking=-48)

left_items = [
    "Pure audio classifier (AASIST, RawNet)",
    "Arms race with cloning tech",
    "Works for researchers.",
    "Useless at the moment of attack.",
]
for i, item in enumerate(left_items):
    yy = col_y + Inches(2.1) + Inches(0.58 * i)
    # bullet dash
    add_text(s, left_x + Inches(0.5), yy, Inches(0.3), Inches(0.4),
             "—", size=14, color=SLATE_GRAY, font=HEAD_FONT)
    add_text(s, left_x + Inches(0.85), yy, col_w - Inches(1.4), Inches(0.5),
             item, size=14, color=SLATE_GRAY if i < 2 else CHARCOAL, font=BODY_FONT)

# Big X stamp
add_text(s, left_x + col_w - Inches(1.3), col_y + Inches(0.3),
         Inches(1.0), Inches(1.0),
         "✕", size=60, bold=True, color=DUST_TAUPE, font=HEAD_FONT,
         align=PP_ALIGN.CENTER)

# RIGHT card — What's missing (accent card)
right_x = left_x + col_w + Inches(0.35)
add_rounded(s, right_x, col_y, col_w, col_h, fill=INK_BLACK, radius_in=0.5)
add_eyebrow(s, right_x + Inches(0.5), col_y + Inches(0.5),
            "• WHAT'S MISSING", color=SIGNAL_LIGHT, text_color=SIGNAL_LIGHT)
add_text(s, right_x + Inches(0.5), col_y + Inches(0.95),
         col_w - Inches(1), Inches(0.8),
         "For everyday users",
         size=24, bold=True, color=CANVAS_CREAM, font=HEAD_FONT, tracking=-48)

right_items = [
    "Real-time, on-phone verification",
    "In Bahasa, English, and Manglish",
    "Something they can do in 30 seconds",
    "Bank freezes happen only AFTER transfer",
]
for i, item in enumerate(right_items):
    yy = col_y + Inches(2.1) + Inches(0.58 * i)
    add_text(s, right_x + Inches(0.5), yy, Inches(0.3), Inches(0.4),
             "→", size=14, bold=True, color=SIGNAL_LIGHT, font=HEAD_FONT)
    add_text(s, right_x + Inches(0.85), yy, col_w - Inches(1.4), Inches(0.5),
             item, size=14, color=CANVAS_CREAM, font=BODY_FONT)

add_footer_mark(s)
add_page_number(s, 4, TOTAL_SLIDES)


# ---------- SLIDE 5 — THE INSIGHT ----------
s = new_slide()
# ghost watermark
add_text(s, Inches(0), Inches(0.1), Inches(13.333), Inches(2),
         "insight",
         size=180, bold=True, color=GHOST_CREAM, font=HEAD_FONT,
         align=PP_ALIGN.CENTER, tracking=-360)

add_eyebrow(s, Inches(5.8), Inches(1.0), "• THE INSIGHT",
            color=SIGNAL_LIGHT, text_color=SLATE_GRAY)

# Statements
add_text(s, Inches(1), Inches(1.8), Inches(11.333), Inches(1.1),
         "We stopped trying to detect the clone.",
         size=40, color=SLATE_GRAY, font=HEAD_FONT, tracking=-80,
         align=PP_ALIGN.CENTER)

add_rich(s, Inches(1), Inches(2.9), Inches(11.333), Inches(1.1),
         [[{"text": "We started helping the user ", "size": 40, "color": INK_BLACK,
            "font": HEAD_FONT, "tracking": -80, "align": PP_ALIGN.CENTER},
           {"text": "outsmart it.", "size": 40, "bold": True, "color": SIGNAL_LIGHT,
            "font": HEAD_FONT, "tracking": -80, "align": PP_ALIGN.CENTER}]],
         align=PP_ALIGN.CENTER)

# Soundwave in middle
add_soundwave(s, Inches(6.666), Inches(4.5), Inches(2.6), bars=11, color=SIGNAL_LIGHT)

# Quote card — centered pill
card_w = Inches(10)
card_x = (prs.slide_width - card_w) / 2
add_rounded(s, card_x, Inches(5.3), card_w, Inches(1.7), fill=INK_BLACK, radius_in=0.5)
add_text(s, card_x, Inches(5.45), card_w, Inches(0.55),
         "Hand them a personalized counter-question.",
         size=22, bold=True, color=CANVAS_CREAM, font=HEAD_FONT,
         align=PP_ALIGN.CENTER, tracking=-44)
add_text(s, card_x, Inches(6.05), card_w, Inches(0.9),
         "Scammers cannot answer questions rooted in private shared memory.\nReal people can.",
         size=14, color=DUST_TAUPE, font=BODY_FONT,
         align=PP_ALIGN.CENTER, line_spacing=1.4)

add_page_number(s, 5, TOTAL_SLIDES)


# ---------- SLIDE 6 — DEMO: UPLOAD ----------
s = new_slide()
add_eyebrow(s, Inches(0.7), Inches(0.55), "• DEMO — STEP 1")
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(1.0),
         "User forwards a suspicious voice note",
         size=34, bold=True, color=INK_BLACK, font=HEAD_FONT, tracking=-68)

# Mock dropzone
mz_x = Inches(0.7)
mz_y = Inches(2.2)
mz_w = Inches(7.2)
mz_h = Inches(4.6)
add_rounded(s, mz_x, mz_y, mz_w, mz_h, fill=LIFTED_CREAM, radius_in=0.55)
# dashed inner border (simulated with colored rounded-rect)
add_rounded(s, mz_x + Inches(0.4), mz_y + Inches(0.4),
            mz_w - Inches(0.8), mz_h - Inches(0.8),
            fill=CANVAS_CREAM, radius_in=0.45,
            line=SIGNAL_LIGHT)

# Big icon (circle)
icon = s.shapes.add_shape(MSO_SHAPE.OVAL,
                          mz_x + mz_w / 2 - Inches(0.6),
                          mz_y + Inches(0.9),
                          Inches(1.2), Inches(1.2))
icon.shadow.inherit = False
icon.fill.solid()
icon.fill.fore_color.rgb = INK_BLACK
icon.line.fill.background()
add_text(s, mz_x + mz_w / 2 - Inches(0.6), mz_y + Inches(0.9),
         Inches(1.2), Inches(1.2),
         "↑", size=44, bold=True, color=SIGNAL_LIGHT, font=HEAD_FONT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, mz_x, mz_y + Inches(2.4), mz_w, Inches(0.6),
         "Drag & drop your voice note",
         size=20, bold=True, color=INK_BLACK, font=HEAD_FONT,
         align=PP_ALIGN.CENTER, tracking=-40)
add_text(s, mz_x, mz_y + Inches(3.0), mz_w, Inches(0.4),
         "or tap to browse  ·  5 seconds of friction",
         size=12, color=SLATE_GRAY, font=BODY_FONT, align=PP_ALIGN.CENTER)

# File chip
add_pill(s, mz_x + Inches(2.1), mz_y + Inches(3.6), Inches(3.0), Inches(0.5), fill=INK_BLACK)
add_text(s, mz_x + Inches(2.1), mz_y + Inches(3.6), Inches(3.0), Inches(0.5),
         "🎧  boss-voicenote.ogg",
         size=11, bold=True, color=CANVAS_CREAM, font=HEAD_FONT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# RIGHT side — metadata inputs & supported
info_x = Inches(8.2)
info_y = Inches(2.2)
info_w = Inches(4.5)

# Caller phone chip
add_eyebrow(s, info_x, info_y, "• OPTIONAL INPUTS", text_color=SLATE_GRAY)
# Phone field mock
add_rounded(s, info_x, info_y + Inches(0.45), info_w, Inches(0.7),
            fill=LIFTED_CREAM, radius_in=0.3, line=INK_BLACK)
add_text(s, info_x + Inches(0.3), info_y + Inches(0.45), info_w, Inches(0.7),
         "+60 12-345 6789",
         size=13, color=INK_BLACK, font=MONO_FONT, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, info_x, info_y + Inches(1.2), info_w, Inches(0.3),
         "Caller phone number (optional)",
         size=9, color=SLATE_GRAY, font=BODY_FONT, tracking=20)

# Role chips
add_text(s, info_x, info_y + Inches(1.7), info_w, Inches(0.3),
         "Claimed relationship",
         size=9, color=SLATE_GRAY, font=BODY_FONT, tracking=20)

roles = [("Boss", True), ("Family", False), ("Friend", False), ("Bank", False)]
rx = info_x
ry = info_y + Inches(2.1)
for name, active in roles:
    chip_w = Inches(0.95 if name != "Family" else 1.05)
    fill = INK_BLACK if active else LIFTED_CREAM
    line = None if active else INK_BLACK
    add_pill(s, rx, ry, chip_w, Inches(0.45), fill=fill, line=line)
    add_text(s, rx, ry, chip_w, Inches(0.45), name,
             size=11, bold=True,
             color=CANVAS_CREAM if active else INK_BLACK, font=HEAD_FONT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rx += chip_w + Inches(0.15)

# Supported block
add_eyebrow(s, info_x, info_y + Inches(3.0), "• SUPPORTED", text_color=SLATE_GRAY)
add_text(s, info_x, info_y + Inches(3.35), info_w, Inches(1.3),
         "WhatsApp · Telegram audio\nBahasa Melayu · English · Manglish\nUp to 20 MB",
         size=12, color=CHARCOAL, font=BODY_FONT, line_spacing=1.6)

add_footer_mark(s)
add_page_number(s, 6, TOTAL_SLIDES)


# ---------- SLIDE 7 — DEMO: AGENTIC STEPPER ----------
s = new_slide()
add_eyebrow(s, Inches(0.7), Inches(0.55), "• DEMO — STEP 2")
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(1.0),
         "Four Gemini agents work in seconds",
         size=34, bold=True, color=INK_BLACK, font=HEAD_FONT, tracking=-68)

# Agent stepper cards — 5 horizontal
agents = [
    ("01", "Transcribe\n& Characterize", "Transcript + voice\nprosody cues", "done"),
    ("02", "Content\nAnalyzer", "Grounded vs. 138\nMalaysian scam phrases", "active"),
    ("03", "Phone\nLookup", "Firestore DB +\nCCID / BNM / MCMC", "pending"),
    ("04", "Challenge\nGenerator", "Crafts 2 personalized\nverification questions", "pending"),
    ("05", "Safety\nPlan", "Verdict + red flags\n+ action plan", "pending"),
]
c_w = Inches(2.38)
c_h = Inches(3.8)
c_gap = Inches(0.1)
total = c_w * 5 + c_gap * 4
start_x = (prs.slide_width - total) / 2
c_y = Inches(2.2)

# Dashed orbit-ish line behind the row
line_y = c_y + c_h / 2
add_rounded(s, start_x + Inches(0.5), line_y, total - Inches(1.0), Inches(0.04),
            fill=DUST_TAUPE, radius_in=0.02)

for i, (num, title, desc, state) in enumerate(agents):
    cx = start_x + i * (c_w + c_gap)
    if state == "active":
        fill = INK_BLACK
        num_color = SIGNAL_LIGHT
        title_color = CANVAS_CREAM
        desc_color = DUST_TAUPE
    elif state == "done":
        fill = LIFTED_CREAM
        num_color = SUCCESS_GREEN
        title_color = INK_BLACK
        desc_color = SLATE_GRAY
    else:
        fill = LIFTED_CREAM
        num_color = DUST_TAUPE
        title_color = SLATE_GRAY
        desc_color = SLATE_GRAY

    add_rounded(s, cx, c_y, c_w, c_h, fill=fill, radius_in=0.4)

    # status pill
    status_label = {"done": "DONE", "active": "RUNNING…", "pending": "PENDING"}[state]
    status_color = {"done": SUCCESS_GREEN, "active": SIGNAL_LIGHT, "pending": DUST_TAUPE}[state]
    add_pill(s, cx + Inches(0.3), c_y + Inches(0.3), Inches(1.0), Inches(0.3),
             fill=status_color)
    add_text(s, cx + Inches(0.3), c_y + Inches(0.3), Inches(1.0), Inches(0.3),
             status_label, size=8, bold=True,
             color=WHITE if state != "pending" else INK_BLACK,
             font=HEAD_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, tracking=40)

    # big number
    add_text(s, cx, c_y + Inches(0.8), c_w, Inches(1.1), num,
             size=54, bold=True, color=num_color, font=HEAD_FONT,
             align=PP_ALIGN.CENTER, tracking=-108)

    # title
    add_text(s, cx + Inches(0.2), c_y + Inches(2.05), c_w - Inches(0.4), Inches(0.85),
             title, size=16, bold=True, color=title_color, font=HEAD_FONT,
             align=PP_ALIGN.CENTER, tracking=-32, line_spacing=1.1)

    # desc
    add_text(s, cx + Inches(0.2), c_y + Inches(2.95), c_w - Inches(0.4), Inches(0.8),
             desc, size=10, color=desc_color, font=BODY_FONT,
             align=PP_ALIGN.CENTER, line_spacing=1.35)

# Bottom note
add_text(s, Inches(0.7), Inches(6.35), Inches(12), Inches(0.5),
         "Live SSE streaming  ·  Transparency panel shows each agent's JSON output.  ·  No black box.",
         size=14, color=CHARCOAL, font=BODY_FONT, align=PP_ALIGN.CENTER, bold=False)

add_footer_mark(s)
add_page_number(s, 7, TOTAL_SLIDES)


# ---------- SLIDE 8 — DEMO: RESULT REPORT ----------
s = new_slide()
add_eyebrow(s, Inches(0.7), Inches(0.55), "• DEMO — STEP 3")
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(1.0),
         "A verdict. And an action.",
         size=34, bold=True, color=INK_BLACK, font=HEAD_FONT, tracking=-68)

# LEFT: verdict card
v_x = Inches(0.7)
v_y = Inches(2.0)
v_w = Inches(5.3)
v_h = Inches(5.0)
add_rounded(s, v_x, v_y, v_w, v_h, fill=LIFTED_CREAM, radius_in=0.55)
add_eyebrow(s, v_x + Inches(0.5), v_y + Inches(0.4), "• VERDICT",
            text_color=SLATE_GRAY)

# HIGH pill
add_pill(s, v_x + Inches(0.5), v_y + Inches(0.8), Inches(2.0), Inches(0.5),
         fill=SIGNAL_ORANGE)
add_text(s, v_x + Inches(0.5), v_y + Inches(0.8), Inches(2.0), Inches(0.5),
         "HIGH SUSPICION",
         size=11, bold=True, color=WHITE, font=HEAD_FONT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, tracking=60)

# Score big
add_text(s, v_x + Inches(0.5), v_y + Inches(1.45), v_w - Inches(1.0), Inches(1.5),
         "80", size=100, bold=True, color=INK_BLACK, font=HEAD_FONT,
         tracking=-200, line_spacing=1.0)
add_text(s, v_x + Inches(2.35), v_y + Inches(2.25), Inches(2.0), Inches(0.5),
         "/ 100", size=20, color=SLATE_GRAY, font=HEAD_FONT, tracking=-40)

# Red flags
add_eyebrow(s, v_x + Inches(0.5), v_y + Inches(3.2), "• RED FLAGS DETECTED",
            color=SIGNAL_ORANGE, text_color=SIGNAL_ORANGE)
flags = ["Urgent tone", "Payment pressure", "Sensitive PIN ask",
         "Authority impersonation", "Deviated phrasing", "Emotional hook"]
fy = v_y + Inches(3.6)
for i, flag in enumerate(flags):
    col = i % 2
    row = i // 2
    fx = v_x + Inches(0.5) + col * Inches(2.3)
    fyy = fy + row * Inches(0.42)
    add_text(s, fx, fyy, Inches(0.25), Inches(0.35), "✕",
             size=11, bold=True, color=SIGNAL_ORANGE, font=HEAD_FONT)
    add_text(s, fx + Inches(0.28), fyy, Inches(2.1), Inches(0.35), flag,
             size=11, color=CHARCOAL, font=BODY_FONT)

# RIGHT: challenge card (hero)
c_x = Inches(6.2)
c_y_r = Inches(2.0)
c_w_r = Inches(6.5)
c_h_r = Inches(5.0)
add_rounded(s, c_x, c_y_r, c_w_r, c_h_r, fill=INK_BLACK, radius_in=0.55)
add_eyebrow(s, c_x + Inches(0.5), c_y_r + Inches(0.5),
            "• THE HERO OUTPUT — COUNTER-QUESTION",
            color=SIGNAL_LIGHT, text_color=SIGNAL_LIGHT)
add_text(s, c_x + Inches(0.5), c_y_r + Inches(0.9), c_w_r - Inches(1.0), Inches(0.55),
         "Ask them this:",
         size=16, color=DUST_TAUPE, font=HEAD_FONT)

# Big quote
add_text(s, c_x + Inches(0.5), c_y_r + Inches(1.5), c_w_r - Inches(1.0), Inches(2.1),
         "“What was the name of the\nrestaurant we went to for\nlunch last Tuesday?”",
         size=24, bold=True, color=CANVAS_CREAM, font=HEAD_FONT,
         tracking=-48, line_spacing=1.25)

# WhatsApp CTA pill
wa_x = c_x + Inches(0.5)
wa_y = c_y_r + Inches(3.7)
add_pill(s, wa_x, wa_y, Inches(3.4), Inches(0.65), fill=SIGNAL_LIGHT)
add_text(s, wa_x, wa_y, Inches(3.4), Inches(0.65),
         "↗  Send via WhatsApp",
         size=14, bold=True, color=INK_BLACK, font=HEAD_FONT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, tracking=-28)

add_text(s, c_x + Inches(0.5), c_y_r + Inches(4.4), c_w_r - Inches(1.0), Inches(0.35),
         "One tap. Opens WhatsApp with the question typed.",
         size=11, color=DUST_TAUPE, font=BODY_FONT, italic=False)

add_footer_mark(s)
add_page_number(s, 8, TOTAL_SLIDES)


# ---------- SLIDE 9 — ARCHITECTURE ----------
s = new_slide()
add_eyebrow(s, Inches(0.7), Inches(0.55), "• ARCHITECTURE")
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(1.0),
         "100% Google Cloud. Region: asia-southeast1.",
         size=30, bold=True, color=INK_BLACK, font=HEAD_FONT, tracking=-60)

# Flow: User → CloudRun(NextJS SSE) → Genkit → Vertex(Gemini 2.5 Flash)
# Firestore below, external portals below right
row_y = Inches(2.15)
box_h = Inches(1.15)

def flow_box(x, y, w, h, label, sub, fill=LIFTED_CREAM, text_color=INK_BLACK, sub_color=SLATE_GRAY):
    add_rounded(s, x, y, w, h, fill=fill, radius_in=0.35)
    add_text(s, x, y + Inches(0.22), w, Inches(0.45), label,
             size=16, bold=True, color=text_color, font=HEAD_FONT,
             align=PP_ALIGN.CENTER, tracking=-32)
    add_text(s, x + Inches(0.15), y + Inches(0.7), w - Inches(0.3), Inches(0.5), sub,
             size=10, color=sub_color, font=BODY_FONT,
             align=PP_ALIGN.CENTER, line_spacing=1.3)

def arrow(x, y, w=Inches(0.5)):
    add_text(s, x, y, w, Inches(0.4), "→",
             size=22, bold=True, color=SIGNAL_LIGHT, font=HEAD_FONT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Row 1
positions_x = [Inches(0.7), Inches(3.4), Inches(6.8), Inches(10.2)]
box_w = Inches(2.4)
boxes1 = [
    ("User", "Mobile / Web\nMalaysia", LIFTED_CREAM, INK_BLACK),
    ("Cloud Run", "Next.js 16\nSSE streaming", INK_BLACK, CANVAS_CREAM),
    ("Firebase Genkit", "4-agent\norchestration", LIFTED_CREAM, INK_BLACK),
    ("Vertex AI", "Gemini 2.5 Flash\nasia-southeast1", INK_BLACK, CANVAS_CREAM),
]
for i, ((label, sub, fill, tc)) in enumerate(boxes1):
    flow_box(positions_x[i], row_y, box_w, box_h, label, sub,
             fill=fill, text_color=tc,
             sub_color=DUST_TAUPE if fill == INK_BLACK else SLATE_GRAY)

# Arrows row 1
arrow_y = row_y + box_h / 2 - Inches(0.2)
arrow(Inches(3.1), arrow_y)
arrow(Inches(6.4), arrow_y)
arrow(Inches(9.85), arrow_y)

# Row 2 — Firestore + external, connected by orbit
row2_y = row_y + box_h + Inches(0.9)
# Firestore center
f_x = Inches(4.4)
add_rounded(s, f_x, row2_y, Inches(4.5), Inches(1.1), fill=LIFTED_CREAM, radius_in=0.35,
            line=SIGNAL_LIGHT)
add_text(s, f_x, row2_y + Inches(0.2), Inches(4.5), Inches(0.4),
         "Firestore (native mode)",
         size=15, bold=True, color=INK_BLACK, font=HEAD_FONT,
         align=PP_ALIGN.CENTER, tracking=-30)
add_text(s, f_x, row2_y + Inches(0.62), Inches(4.5), Inches(0.5),
         "scam_phrases (138 BM/EN/Manglish)  ·  scam_numbers",
         size=10, color=SLATE_GRAY, font=BODY_FONT, align=PP_ALIGN.CENTER)

# Arrow from CloudRun to Firestore (visual hint)
add_text(s, Inches(3.6), row_y + box_h + Inches(0.05), Inches(2.0), Inches(0.4),
         "↕  Firestore lookups",
         size=9, bold=True, color=SIGNAL_LIGHT, font=HEAD_FONT,
         align=PP_ALIGN.CENTER, tracking=20)

# Row 3 external portals
row3_y = row2_y + Inches(1.1) + Inches(0.55)
portals = [
    ("Semakmule", "CCID police DB"),
    ("BNM FCAL", "Financial Consumer Alert List"),
    ("MCMC", "Official telco sources"),
]
p_w = Inches(3.8)
p_gap = Inches(0.25)
p_total = p_w * 3 + p_gap * 2
p_start = (prs.slide_width - p_total) / 2
for i, (label, sub) in enumerate(portals):
    px = p_start + i * (p_w + p_gap)
    add_pill(s, px, row3_y, p_w, Inches(0.85), fill=CANVAS_CREAM, line=INK_BLACK)
    add_text(s, px, row3_y + Inches(0.12), p_w, Inches(0.4), label,
             size=13, bold=True, color=INK_BLACK, font=HEAD_FONT,
             align=PP_ALIGN.CENTER, tracking=-26)
    add_text(s, px, row3_y + Inches(0.45), p_w, Inches(0.4), sub,
             size=9, color=SLATE_GRAY, font=BODY_FONT,
             align=PP_ALIGN.CENTER, line_spacing=1.2)

# External label
add_eyebrow(s, p_start, row3_y - Inches(0.4), "• EXTERNAL DEEP LINKS")

# Bottom tag line
add_text(s, Inches(0.7), Inches(7.1), Inches(12), Inches(0.3),
         "~50ms latency from Malaysia  ·  No API keys in production  ·  IAM via Cloud Run default SA",
         size=10, color=SLATE_GRAY, font=HEAD_FONT, align=PP_ALIGN.CENTER, tracking=20)

add_page_number(s, 9, TOTAL_SLIDES)


# ---------- SLIDE 10 — STACK: 100% GOOGLE AI ----------
s = new_slide()
add_eyebrow(s, Inches(0.7), Inches(0.55), "• STACK")
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(1.0),
         "Google AI ecosystem, end-to-end.",
         size=34, bold=True, color=INK_BLACK, font=HEAD_FONT, tracking=-68)

stack = [
    ("MODEL", "Gemini 2.5 Flash", "Pro escalation available", SIGNAL_LIGHT),
    ("RUNTIME", "Vertex AI", "asia-southeast1 region", INK_BLACK),
    ("ORCHESTRATION", "Firebase Genkit JS", "Typed, streaming, inspectable", SIGNAL_LIGHT),
    ("DATA", "Firestore", "Native mode, 138-phrase corpus", INK_BLACK),
    ("DEPLOY", "Google Cloud Run", "Serverless · SSE · auto-scale", INK_BLACK),
    ("PROMPT ITERATION", "AI Studio", "Grounded prompt tuning", SIGNAL_LIGHT),
    ("DEV", "Antigravity + Genkit Dev UI", "Local agent flow debugger", INK_BLACK),
]

# 4-col grid on top row, 3 on bottom — like an editorial magazine grid
cols = 4
card_w = Inches(2.95)
card_h = Inches(1.9)
gap_x = Inches(0.15)
gap_y = Inches(0.25)
grid_w = card_w * cols + gap_x * (cols - 1)
grid_x = (prs.slide_width - grid_w) / 2
grid_y = Inches(2.2)

for i, (eyebrow, name, sub, accent) in enumerate(stack):
    row = i // cols
    col = i % cols
    x = grid_x + col * (card_w + gap_x)
    y = grid_y + row * (card_h + gap_y)

    fill = LIFTED_CREAM if accent != INK_BLACK else INK_BLACK
    is_dark = fill == INK_BLACK
    add_rounded(s, x, y, card_w, card_h, fill=fill, radius_in=0.35)

    # eyebrow
    dot_color = SIGNAL_LIGHT
    add_eyebrow(s, x + Inches(0.35), y + Inches(0.3), f"• {eyebrow}",
                color=dot_color,
                text_color=DUST_TAUPE if is_dark else SLATE_GRAY)
    # name
    add_text(s, x + Inches(0.35), y + Inches(0.7), card_w - Inches(0.7), Inches(0.6),
             name, size=18, bold=True,
             color=CANVAS_CREAM if is_dark else INK_BLACK,
             font=HEAD_FONT, tracking=-36)
    # sub
    add_text(s, x + Inches(0.35), y + Inches(1.25), card_w - Inches(0.7), Inches(0.6),
             sub, size=10,
             color=DUST_TAUPE if is_dark else SLATE_GRAY,
             font=BODY_FONT, line_spacing=1.3)

# Second row only has 3; draw a decorative caption in the empty 4th slot
empty_x = grid_x + 3 * (card_w + gap_x)
empty_y = grid_y + (card_h + gap_y)
add_text(s, empty_x + Inches(0.2), empty_y + Inches(0.5),
         card_w - Inches(0.4), Inches(1.4),
         "Every\nlayer.\nGoogle.",
         size=28, bold=True, color=SIGNAL_LIGHT,
         font=HEAD_FONT, tracking=-56, line_spacing=1.0,
         align=PP_ALIGN.LEFT)

add_footer_mark(s)
add_page_number(s, 10, TOTAL_SLIDES)


# ---------- SLIDE 11 — WHY AGENTIC ----------
s = new_slide()
add_eyebrow(s, Inches(0.7), Inches(0.55), "• WHY 4 AGENTS, NOT 1 PROMPT")
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(1.0),
         "Focus. Structure. Transparency.",
         size=34, bold=True, color=INK_BLACK, font=HEAD_FONT, tracking=-68)

# 4 agent mini-cards, each with tiny JSON shape
agents_row = [
    ("Transcribe",  "audio",
        '{\n  "transcript": …,\n  "voice_observations": […]\n}'),
    ("Content",     "transcript + RAG",
        '{\n  "urgency_score": 80,\n  "red_flags": […],\n  "sensitive_requests": […]\n}'),
    ("Challenge",   "transcript + content",
        '{\n  "questions": […],\n  "whatsapp_copypaste": …\n}'),
    ("Safety Plan", "everything",
        '{\n  "verdict": "HIGH",\n  "suspicion_score": 80,\n  "action_plan": […]\n}'),
]
col_w = Inches(2.95)
col_h = Inches(3.6)
gx = Inches(0.15)
tw = col_w * 4 + gx * 3
sx = (prs.slide_width - tw) / 2
sy = Inches(2.2)

for i, (name, inp, out) in enumerate(agents_row):
    x = sx + i * (col_w + gx)
    add_rounded(s, x, sy, col_w, col_h, fill=LIFTED_CREAM, radius_in=0.35)
    # header pill
    add_pill(s, x + Inches(0.3), sy + Inches(0.3), Inches(0.8), Inches(0.3),
             fill=INK_BLACK)
    add_text(s, x + Inches(0.3), sy + Inches(0.3), Inches(0.8), Inches(0.3),
             f"0{i+1}", size=10, bold=True, color=SIGNAL_LIGHT, font=HEAD_FONT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, tracking=40)
    # title
    add_text(s, x + Inches(0.3), sy + Inches(0.75), col_w - Inches(0.6), Inches(0.5),
             name, size=20, bold=True, color=INK_BLACK, font=HEAD_FONT, tracking=-40)
    # input row
    add_text(s, x + Inches(0.3), sy + Inches(1.3), col_w - Inches(0.6), Inches(0.3),
             "INPUT", size=8, bold=True, color=SLATE_GRAY, font=HEAD_FONT, tracking=32)
    add_text(s, x + Inches(0.3), sy + Inches(1.55), col_w - Inches(0.6), Inches(0.3),
             inp, size=11, color=CHARCOAL, font=BODY_FONT, italic=True)
    # output JSON
    add_text(s, x + Inches(0.3), sy + Inches(1.95), col_w - Inches(0.6), Inches(0.3),
             "OUTPUT (Zod-enforced)", size=8, bold=True, color=SLATE_GRAY, font=HEAD_FONT,
             tracking=32)
    add_rounded(s, x + Inches(0.3), sy + Inches(2.2), col_w - Inches(0.6), Inches(1.2),
                fill=INK_BLACK, radius_in=0.15)
    add_text(s, x + Inches(0.45), sy + Inches(2.3), col_w - Inches(0.9), Inches(1.0),
             out, size=9, color=SIGNAL_LIGHT, font=MONO_FONT, line_spacing=1.3)

# Three benefits pills at bottom
benefits = [
    ("Cheaper", "ChallengeGen never sees audio"),
    ("Transparent", "Every agent inspectable in UI"),
    ("Robust", "Zod rejects bad model output"),
]
by = Inches(6.1)
b_w = Inches(3.8)
b_gap = Inches(0.25)
b_total = b_w * 3 + b_gap * 2
b_x = (prs.slide_width - b_total) / 2
for i, (t, d) in enumerate(benefits):
    bx = b_x + i * (b_w + b_gap)
    add_pill(s, bx, by, b_w, Inches(0.7), fill=INK_BLACK)
    add_rich(s, bx + Inches(0.2), by + Inches(0.1), b_w - Inches(0.4), Inches(0.55),
             [[{"text": t + "  ", "size": 13, "bold": True, "color": SIGNAL_LIGHT,
                "font": HEAD_FONT, "tracking": -26, "align": PP_ALIGN.CENTER},
               {"text": "— " + d, "size": 12, "color": CANVAS_CREAM,
                "font": BODY_FONT, "align": PP_ALIGN.CENTER}]])

add_page_number(s, 11, TOTAL_SLIDES)


# ---------- SLIDE 12 — MALAYSIAN IMPACT ----------
s = new_slide()
add_eyebrow(s, Inches(0.7), Inches(0.55), "• BUILT FOR MALAYSIA")
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(1.0),
         "Not a translation. A localization.",
         size=34, bold=True, color=INK_BLACK, font=HEAD_FONT, tracking=-68)

# Three content columns
boxes = [
    ("BILINGUAL + MANGLISH",
     "138 phrases",
     "59 BM  ·  57 EN\n22 Manglish",
     "The voice note that says\n“cepat lah wei, urgent ni”\ngets understood.",
     SIGNAL_LIGHT),
    ("LOCAL HOTLINES",
     "Direct tel:",
     "CCID 03-2610 1559\nBNM 1-300-88-5465\n999 Emergency",
     "Built into every\nsafety plan.",
     INK_BLACK),
    ("AUTHORITATIVE",
     "Semakmule\nBNM  ·  MCMC",
     "Official deep links",
     "We point users to the\nauthority — we don't\nscrape it.",
     SIGNAL_LIGHT),
]

col_w = Inches(4.0)
col_h = Inches(4.6)
gap = Inches(0.25)
total_w = col_w * 3 + gap * 2
start_x = (prs.slide_width - total_w) / 2
col_y = Inches(2.2)

for i, (eyebrow, headline, sub, body, accent) in enumerate(boxes):
    x = start_x + i * (col_w + gap)
    fill = INK_BLACK if i == 1 else LIFTED_CREAM
    is_dark = fill == INK_BLACK
    add_rounded(s, x, col_y, col_w, col_h, fill=fill, radius_in=0.5)

    add_eyebrow(s, x + Inches(0.45), col_y + Inches(0.5),
                f"• {eyebrow}", color=accent,
                text_color=DUST_TAUPE if is_dark else SLATE_GRAY)

    add_text(s, x + Inches(0.45), col_y + Inches(1.05),
             col_w - Inches(0.9), Inches(1.3),
             headline, size=30, bold=True,
             color=CANVAS_CREAM if is_dark else INK_BLACK,
             font=HEAD_FONT, tracking=-60, line_spacing=1.05)

    add_text(s, x + Inches(0.45), col_y + Inches(2.5),
             col_w - Inches(0.9), Inches(0.9),
             sub, size=12, bold=True,
             color=accent if not is_dark else SIGNAL_LIGHT,
             font=HEAD_FONT, line_spacing=1.35)

    add_text(s, x + Inches(0.45), col_y + Inches(3.45),
             col_w - Inches(0.9), Inches(1.1),
             body, size=12,
             color=DUST_TAUPE if is_dark else CHARCOAL,
             font=BODY_FONT, italic=True, line_spacing=1.4)

add_footer_mark(s)
add_page_number(s, 12, TOTAL_SLIDES)


# ---------- SLIDE 13 — TRACTION PATH ----------
s = new_slide()
add_eyebrow(s, Inches(0.7), Inches(0.55), "• PATH TO ADOPTION")
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(1.0),
         "From live app → national partnerships.",
         size=34, bold=True, color=INK_BLACK, font=HEAD_FONT, tracking=-68)

phases = [
    ("NOW", "Phase 1", "Public web app",
     ["Shareable via WhatsApp", "Free. No login.", "You can test it today."],
     SIGNAL_LIGHT),
    ("NEXT 6 MONTHS", "Phase 2", "Scale & partner",
     ["WhatsApp Business webhook",
      "Partner with PDRM CCID & BNM",
      "Schools: cyber-safety"],
     INK_BLACK),
    ("12 MONTHS", "Phase 3", "Deeper product",
     ["Mobile app + browser extension",
      "Deepfake classifier signal",
      "Community phone database"],
     CANVAS_CREAM),
]

col_w = Inches(4.0)
col_h = Inches(4.7)
gap = Inches(0.25)
total_w = col_w * 3 + gap * 2
start_x = (prs.slide_width - total_w) / 2
col_y = Inches(2.15)

# Timeline line behind
line_y = col_y + Inches(0.9)
add_rounded(s, start_x + Inches(0.5), line_y,
            total_w - Inches(1.0), Inches(0.04),
            fill=DUST_TAUPE, radius_in=0.02)

for i, (eyebrow, phase, title, bullets, accent) in enumerate(phases):
    x = start_x + i * (col_w + gap)
    is_active = (i == 0)
    fill = INK_BLACK if is_active else LIFTED_CREAM
    is_dark = fill == INK_BLACK
    add_rounded(s, x, col_y, col_w, col_h, fill=fill, radius_in=0.5)

    # timeline dot
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                             x + col_w / 2 - Inches(0.15),
                             line_y - Inches(0.125),
                             Inches(0.3), Inches(0.3))
    dot.shadow.inherit = False
    dot.fill.solid()
    dot.fill.fore_color.rgb = SIGNAL_LIGHT if is_active else DUST_TAUPE
    dot.line.color.rgb = CANVAS_CREAM
    dot.line.width = Pt(3)

    add_eyebrow(s, x + Inches(0.45), col_y + Inches(0.5),
                f"• {eyebrow}",
                color=SIGNAL_LIGHT if is_active else DUST_TAUPE,
                text_color=DUST_TAUPE if is_dark else SLATE_GRAY)

    # phase label
    add_text(s, x + Inches(0.45), col_y + Inches(1.4),
             col_w - Inches(0.9), Inches(0.5),
             phase, size=14, color=SIGNAL_LIGHT if is_dark else SIGNAL_LIGHT,
             font=HEAD_FONT, bold=True, tracking=-28)

    # title
    add_text(s, x + Inches(0.45), col_y + Inches(1.8),
             col_w - Inches(0.9), Inches(1.0),
             title, size=26, bold=True,
             color=CANVAS_CREAM if is_dark else INK_BLACK,
             font=HEAD_FONT, tracking=-52, line_spacing=1.1)

    # bullets
    for j, b in enumerate(bullets):
        by = col_y + Inches(2.95) + Inches(0.55) * j
        add_text(s, x + Inches(0.45), by, Inches(0.25), Inches(0.4),
                 "→", size=12, bold=True, color=SIGNAL_LIGHT, font=HEAD_FONT)
        add_text(s, x + Inches(0.75), by,
                 col_w - Inches(1.2), Inches(0.5),
                 b, size=12,
                 color=CANVAS_CREAM if is_dark else CHARCOAL,
                 font=BODY_FONT, line_spacing=1.3)

add_footer_mark(s)
add_page_number(s, 13, TOTAL_SLIDES)


# ---------- SLIDE 14 — ROADMAP + HONEST LIMITS ----------
s = new_slide()
add_eyebrow(s, Inches(0.7), Inches(0.55), "• ROADMAP & HONEST LIMITS")
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(1.0),
         "What's next — and what we don't claim.",
         size=32, bold=True, color=INK_BLACK, font=HEAD_FONT, tracking=-64)

col_w = Inches(5.9)
col_h = Inches(4.0)
col_y = Inches(2.2)
gap = Inches(0.35)

# LEFT — Next iterations (cream + green tint)
lx = Inches(0.7)
add_rounded(s, lx, col_y, col_w, col_h, fill=LIFTED_CREAM, radius_in=0.55)
add_eyebrow(s, lx + Inches(0.5), col_y + Inches(0.5),
            "• WHAT'S NEXT", color=SUCCESS_GREEN, text_color=SUCCESS_GREEN)
add_text(s, lx + Inches(0.5), col_y + Inches(0.95),
         col_w - Inches(1), Inches(0.6),
         "Next iterations",
         size=22, bold=True, color=INK_BLACK, font=HEAD_FONT, tracking=-44)

next_items = [
    "Audio > 20 MB via Files API",
    "Gemini TTS playback of the challenge question",
    "Firebase Auth + per-user history",
    "Vertex AI Vector Search (when corpus > 1k items)",
]
for i, it in enumerate(next_items):
    yy = col_y + Inches(1.75) + Inches(0.5) * i
    add_text(s, lx + Inches(0.5), yy, Inches(0.3), Inches(0.5),
             "✓", size=14, bold=True, color=SUCCESS_GREEN, font=HEAD_FONT)
    add_text(s, lx + Inches(0.85), yy, col_w - Inches(1.35), Inches(0.5),
             it, size=13, color=CHARCOAL, font=BODY_FONT)

# RIGHT — honest limits (ink card)
rx = lx + col_w + gap
add_rounded(s, rx, col_y, col_w, col_h, fill=INK_BLACK, radius_in=0.55)
add_eyebrow(s, rx + Inches(0.5), col_y + Inches(0.5),
            "• WE ARE HONEST THAT", color=SIGNAL_LIGHT, text_color=SIGNAL_LIGHT)
add_text(s, rx + Inches(0.5), col_y + Inches(0.95),
         col_w - Inches(1), Inches(0.6),
         "What we don't claim",
         size=22, bold=True, color=CANVAS_CREAM, font=HEAD_FONT, tracking=-44)

honest_items = [
    "This is NOT a binary deepfake confirmation",
    "Voice cues alone cannot prove synthesis",
    "Our phone DB is curated (30 numbers), not exhaustive",
    "For authoritative checks we link to CCID / BNM / MCMC",
]
for i, it in enumerate(honest_items):
    yy = col_y + Inches(1.75) + Inches(0.5) * i
    add_text(s, rx + Inches(0.5), yy, Inches(0.3), Inches(0.5),
             "—", size=14, bold=True, color=SIGNAL_LIGHT, font=HEAD_FONT)
    add_text(s, rx + Inches(0.85), yy, col_w - Inches(1.35), Inches(0.5),
             it, size=13, color=CANVAS_CREAM, font=BODY_FONT)

# Bottom callout pill — posture
posture_y = Inches(6.45)
pw = Inches(12)
px = (prs.slide_width - pw) / 2
add_pill(s, px, posture_y, pw, Inches(0.65), fill=CANVAS_CREAM, line=INK_BLACK)
add_rich(s, px + Inches(0.3), posture_y + Inches(0.1), pw - Inches(0.6), Inches(0.5),
         [[{"text": "POSTURE  ·  ", "size": 12, "bold": True, "color": SIGNAL_LIGHT,
            "font": HEAD_FONT, "tracking": 60, "align": PP_ALIGN.CENTER},
           {"text": "suspicion score", "size": 14, "bold": True, "color": INK_BLACK,
            "font": HEAD_FONT, "italic": True, "align": PP_ALIGN.CENTER},
           {"text": ", not deepfake verdict. This survives adversarial Q&A.",
            "size": 14, "color": CHARCOAL, "font": BODY_FONT,
            "align": PP_ALIGN.CENTER}]])

add_page_number(s, 14, TOTAL_SLIDES)


# ---------- SLIDE 15 — CLOSING · LIVE URL · QR ----------
s = new_slide()

# ghost watermark at top
add_text(s, Inches(0), Inches(0.1), Inches(13.333), Inches(2),
         "try it",
         size=180, bold=True, color=GHOST_CREAM, font=HEAD_FONT,
         align=PP_ALIGN.CENTER, tracking=-360)

add_eyebrow(s, Inches(6.0), Inches(1.15), "• TRY IT NOW",
            color=SIGNAL_LIGHT, text_color=SLATE_GRAY)

# QR placeholder — big rounded square on left
qr_x = Inches(1.5)
qr_y = Inches(1.9)
qr_size = Inches(4.2)
add_rounded(s, qr_x, qr_y, qr_size, qr_size, fill=LIFTED_CREAM, radius_in=0.4,
            line=INK_BLACK)

# Fake QR pattern made of small squares
import random
random.seed(42)
cells = 21
cell_size = Emu(qr_size / cells)
pad = Emu(qr_size / 12)  # margin inside card
inner_x = qr_x + pad
inner_y = qr_y + pad
inner_size = qr_size - 2 * pad
cell_size = Emu(inner_size / cells)
for cy in range(cells):
    for cx in range(cells):
        # Create finder patterns in corners
        in_tl = cx < 7 and cy < 7
        in_tr = cx >= cells - 7 and cy < 7
        in_bl = cx < 7 and cy >= cells - 7
        in_finder = in_tl or in_tr or in_bl
        if in_finder:
            # draw full square if outer ring or inner 3x3
            if in_tl:
                ax, ay = cx, cy
            elif in_tr:
                ax, ay = cx - (cells - 7), cy
            else:
                ax, ay = cx, cy - (cells - 7)
            # outer border
            if ax in (0, 6) or ay in (0, 6):
                draw = True
            elif 2 <= ax <= 4 and 2 <= ay <= 4:
                draw = True
            else:
                draw = False
        else:
            draw = random.random() < 0.45
        if draw:
            bx = inner_x + cx * cell_size
            by = inner_y + cy * cell_size
            sq = slide_bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, by, cell_size, cell_size)
            sq.shadow.inherit = False
            sq.fill.solid()
            sq.fill.fore_color.rgb = INK_BLACK
            sq.line.fill.background()

# QR caption
add_text(s, qr_x, qr_y + qr_size + Inches(0.15), qr_size, Inches(0.3),
         "SCAN TO TRY  ·  asia-southeast1.run.app",
         size=9, bold=True, color=SLATE_GRAY, font=HEAD_FONT,
         align=PP_ALIGN.CENTER, tracking=40)

# RIGHT — URL, tagline, team
rx = Inches(6.3)
add_text(s, rx, Inches(1.9), Inches(6.5), Inches(0.4),
         "🌐 LIVE URL",
         size=11, bold=True, color=SIGNAL_LIGHT, font=HEAD_FONT,
         tracking=44)
add_text(s, rx, Inches(2.35), Inches(6.5), Inches(1.4),
         "dengardulu-169906713421\n.asia-southeast1.run.app",
         size=18, bold=True, color=INK_BLACK, font=MONO_FONT,
         line_spacing=1.2, tracking=-20)

# Tagline (BM + EN stacked)
add_text(s, rx, Inches(4.15), Inches(6.8), Inches(0.7),
         "Dengar dulu. Jawab kemudian.",
         size=26, bold=True, color=INK_BLACK, font=HEAD_FONT, tracking=-52)
add_text(s, rx, Inches(4.85), Inches(6.8), Inches(0.55),
         "Listen first. Answer wisely.",
         size=19, color=SIGNAL_LIGHT, font=HEAD_FONT, tracking=-38, italic=True)

# Team pill
team_y = Inches(5.85)
add_pill(s, rx, team_y, Inches(6.3), Inches(0.55), fill=INK_BLACK)
add_text(s, rx, team_y, Inches(6.3), Inches(0.55),
         "— TEAM DENGARDULU  ·  GDG ON CAMPUS UTM  ·  APR 2026 —",
         size=10, bold=True, color=CANVAS_CREAM, font=HEAD_FONT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, tracking=60)

# Thank you soundwave
add_soundwave(s, Inches(9.45), Inches(6.8), Inches(2.0), bars=9, color=SIGNAL_LIGHT)
add_text(s, rx, Inches(6.75), Inches(3.0), Inches(0.3),
         "thank you.",
         size=13, color=SLATE_GRAY, font=HEAD_FONT,
         italic=True, tracking=-26)

add_page_number(s, 15, TOTAL_SLIDES)


# ---------- Save ----------
import os
out_path = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "pitch-deck.pptx",
)
prs.save(out_path)
print(f"OK → {out_path}")
print(f"Slides: {len(prs.slides)}")
print(f"Size: {prs.slide_width.inches}\" x {prs.slide_height.inches}\"")
